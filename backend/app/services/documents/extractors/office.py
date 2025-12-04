import io
import logging
from typing import List, Optional

import pandas as pd
from charset_normalizer import from_bytes
from docx import Document
from pptx import Presentation
from striprtf.striprtf import rtf_to_text

from app.services.documents.exceptions import ExtractionError
from app.services.documents.extractors.base import BaseExtractor
from app.services.documents.types import DocumentContent, ExtractionOptions

logger = logging.getLogger(__name__)


class DocxExtractor(BaseExtractor):
    canonical_types: set[str] = {"docx"}
    mime_types: set[str] = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    }
    extensions: set[str] = {"docx", "doc"}
    primary_type: str = "docx"

    def extract(self, data: bytes, filename: str | None, content_type: str | None, options: ExtractionOptions) -> DocumentContent:
        warnings: list[str] = []
        try:
            doc = Document(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
            tables = len(doc.tables)
            metadata = {"paragraphs": len(paragraphs), "tables": tables}
            text = "\n\n".join(paragraphs)
        except Exception as exc:
            raise ExtractionError(f"Unable to parse DOCX: {exc}") from exc

        text = self.trim_text(self.normalize_whitespace(text), warnings)
        return DocumentContent(
            text=text,
            metadata=metadata,
            file_type="docx",
            content_type=content_type or "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            warnings=warnings,
        )


class RtfExtractor(BaseExtractor):
    canonical_types: set[str] = {"rtf"}
    mime_types: set[str] = {"application/rtf", "text/rtf"}
    extensions: set[str] = {"rtf"}
    primary_type: str = "rtf"

    def extract(self, data: bytes, filename: str | None, content_type: str | None, options: ExtractionOptions) -> DocumentContent:
        warnings: list[str] = []
        try:
            text = rtf_to_text(data.decode("utf-8", errors="ignore"))
        except Exception as exc:
            raise ExtractionError(f"Unable to parse RTF: {exc}") from exc

        text = self.trim_text(self.normalize_whitespace(text), warnings)
        return DocumentContent(
            text=text,
            metadata={"source": filename or "upload"},
            file_type="rtf",
            content_type=content_type or "text/rtf",
            warnings=warnings,
        )


class PptxExtractor(BaseExtractor):
    canonical_types: set[str] = {"pptx"}
    mime_types: set[str] = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }
    extensions: set[str] = {"pptx", "ppt"}
    primary_type: str = "pptx"

    def extract(self, data: bytes, filename: str | None, content_type: str | None, options: ExtractionOptions) -> DocumentContent:
        warnings: list[str] = []
        try:
            presentation = Presentation(io.BytesIO(data))
            slide_texts: List[str] = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                chunks: List[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:  # pragma: no cover - optional path
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text:
                        chunks.append(f"Notes: {notes_text}")
                if chunks:
                    slide_texts.append(f"--- Slide {slide_index} ---\n" + "\n".join(chunks))

            text = "\n\n".join(slide_texts)
            metadata = {"slides": len(presentation.slides)}
        except Exception as exc:
            raise ExtractionError(f"Unable to parse PPTX: {exc}") from exc

        text = self.trim_text(self.normalize_whitespace(text), warnings)
        return DocumentContent(
            text=text,
            metadata=metadata,
            file_type="pptx",
            content_type=content_type or "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            warnings=warnings,
        )


class SpreadsheetExtractor(BaseExtractor):
    canonical_types: set[str] = {"xlsx", "xls", "ods"}
    mime_types: set[str] = {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
        "application/vnd.oasis.opendocument.spreadsheet",
    }
    extensions: set[str] = {"xlsx", "xls", "ods"}
    primary_type: str = "xlsx"

    def __init__(self, row_preview: int = 20):
        self.row_preview = row_preview

    def extract(self, data: bytes, filename: str | None, content_type: str | None, options: ExtractionOptions) -> DocumentContent:
        warnings: list[str] = []
        try:
            excel_file = pd.ExcelFile(io.BytesIO(data))
        except Exception as exc:
            raise ExtractionError(f"Unable to open spreadsheet: {exc}") from exc

        text_blocks: List[str] = []
        metadata: dict[str, int | str] = {"sheets": len(excel_file.sheet_names)}

        for sheet_name in excel_file.sheet_names:
            try:
                df = pd.read_excel(io.BytesIO(data), sheet_name=sheet_name, engine=self._engine_for_sheet(filename))
            except Exception as exc:  # pragma: no cover - engine edge cases
                warnings.append(f"Failed to read sheet {sheet_name}: {exc}")
                continue

            display_df = df.head(self.row_preview)
            sheet_header = f"--- Sheet: {sheet_name} ({len(df)} rows, {len(df.columns)} columns) ---"
            text_blocks.append(sheet_header)
            text_blocks.append(display_df.to_string(index=False))
            if len(df) > self.row_preview:
                tail_preview = df.tail(3)
                text_blocks.append("... (truncated) ...")
                text_blocks.append(tail_preview.to_string(index=False))
                warnings.append(f"Sheet {sheet_name} truncated to first {self.row_preview} rows")

        text = "\n\n".join(text_blocks)
        text = self.trim_text(self.normalize_whitespace(text), warnings)

        return DocumentContent(
            text=text,
            metadata=metadata,
            file_type="xlsx",
            content_type=content_type or "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            warnings=warnings,
        )

    def _engine_for_sheet(self, filename: Optional[str]) -> Optional[str]:
        if not filename:
            return None
        lowered = filename.lower()
        if lowered.endswith(".ods"):
            return "odf"
        return None


class CsvExtractor(BaseExtractor):
    canonical_types: set[str] = {"csv"}
    mime_types: set[str] = {"text/csv"}
    extensions: set[str] = {"csv"}
    primary_type: str = "csv"

    def __init__(self, row_preview: int = 200):
        self.row_preview = row_preview

    def extract(self, data: bytes, filename: str | None, content_type: str | None, options: ExtractionOptions) -> DocumentContent:
        warnings: list[str] = []
        encoding = self._detect_encoding(data)
        try:
            df = pd.read_csv(io.BytesIO(data), nrows=self.row_preview, encoding=encoding, encoding_errors="ignore")
        except Exception as exc:
            raise ExtractionError(f"Unable to parse CSV: {exc}") from exc

        total_rows = len(df)
        if total_rows >= self.row_preview:
            warnings.append(f"CSV truncated to first {self.row_preview} rows for preview")

        metadata = {
            "rows": total_rows,
            "columns": len(df.columns),
            "encoding": encoding or "utf-8",
        }

        text = df.to_string(index=False)
        text = self.trim_text(self.normalize_whitespace(text), warnings)
        return DocumentContent(
            text=text,
            metadata=metadata,
            file_type="csv",
            content_type=content_type or "text/csv",
            warnings=warnings,
        )

    def _detect_encoding(self, data: bytes) -> Optional[str]:
        result = from_bytes(data).best()
        return result.encoding if result else None
